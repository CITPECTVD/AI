from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define some generic slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    blank_slide_layout = prs.slide_layouts[6]
    
    # helper for adding title and text
    def add_bullet_slide(title, points):
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        title_shape.text = title
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
                p.text = point
            else:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0
            p.font.size = Pt(20)

    # 1. Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "AI 影片製程與生成簡報"
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(35, 86, 157) # primary blue
    subtitle.text = "TPECT China Airlines\n2026 KEfilm. AI-Enhanced Video Production"

    # 2. Theme & Concept
    add_bullet_slide(
        "1. 影片主題 & 2. 核心概念",
        [
            "【影片主題】：「春節團圓」與「海外旅遊」的聯結。",
            " 展現了農曆新年的傳統習俗，並透過「旅行」賦予團圓新的意義。",
            "",
            "【核心概念】：「有你的旅程，哪裡都是家」。",
            " 品牌強調旅行不只是離開家，而是與心愛的人一起創造共同的回憶，將家的感覺延伸到世界的任何角落。"
        ]
    )

    # 3. Visual Style
    add_bullet_slide(
        "3. 視覺風格",
        [
            "【風格定位】: 真人風格 (Live-action)。寫實的人物與居家場景，營造強烈的親近感與生活氣息。",
            "【色調】: 溫暖且飽和。以企業色「紅色」為視覺主軸，搭配黃光，最後轉向夕陽餘暉的粉紫色澤。",
            "【運鏡】: 中景與特寫交替。捕捉家人笑容、Surprise眼神，及登機證等關鍵物件細節。"
        ]
    )

    # 4. Plot Details
    add_bullet_slide(
        "4. 劇情概要 & 5. 劇情細節",
        [
            "【概要】: 兒子一家三口回到奶奶家過年。全家人的年節活動後，在紅包環節發現驚喜，準備出國旅遊。",
            "",
            "【關鍵細節】:",
            " - 傳統習俗：寫春聯、貼「囍」字貼紙。",
            " - 圍爐年菜：吃火鍋，展示肉片「花開富貴」與「年年有餘」的魚。",
            " - 家庭娛樂：玩大富翁型態桌遊「環遊世界」。",
            " - 拜年驚喜：紅包中裝著前往東京（TOKYO）的登機證。",
            " - 開心出發：穿上紅色系服裝、收拾行李步出家門。"
        ]
    )

    # 5. Characters
    add_bullet_slide(
        "6. 人物形象",
        [
            "【奶奶】 (約 60-70 歲)",
            " 慈祥、充滿驚喜，凝聚家族情感的核心。穿著典雅的紅白色系拼接上衣，配戴珍珠項鍊。",
            "【兒子】 (約 30-35 歲)",
            " 孝順且充滿活力的現代形象。穿著米白色高領毛衣，營造溫暖、可靠氣質。",
            "【媳婦】 (約 30-35 歲)",
            " 孝順且充滿活力的現代形象。穿著亮紅色的露肩針織上衣，展現現代時尚感。",
            "【小女孩】 (約 5-7 歲)",
            " 天真無邪，展現對旅程的期待。綁雙丸子頭，穿著米色針織開襟衫搭配粉色系格子裙。"
        ]
    )

    # 6. Props and Music
    add_bullet_slide(
        "7. 其他形象 (場景/道具) & 8. 音樂型態",
        [
            "【場景/道具】",
            " - 關鍵場景：喜氣洋洋的奶奶家客廳、溫馨圍爐餐桌、準備出發的居家門口。",
            " - 核心道具：中華航空登機證(東京)、大富翁桌遊「環遊世界」、「花開富貴」肉片年菜、春聯。",
            " - 氛圍元素：紅包袋、華航企業色紅白系服裝、室內溫暖黃光、夕陽粉紫色澤。",
            "",
            "【音樂型態】",
            " - 風格：輕快、和諧且節奏明亮的配樂。",
            " - 層次：隨劇情發展由溫馨平穩轉向高昂活潑，呼應及強化全家出國旅行的興奮心情。"
        ]
    )

    # 7. AI Budget
    add_bullet_slide(
        "影片生成費用概括",
        [
            "【圖片生成費用】",
            " - 假設 5s 一張圖，300s 共 60 張，每張 15 次，再 * 2 = 1800 張",
            " - 總價約 $7,718 ~ $13,824 TWD",
            "",
            "【影片生成費用 (1080p)】",
            " - Google Veo (8 秒影片)：單次較便宜，約 $3,645 ~ $18,225 TWD 加上額外另購。",
            " - Kling (5 秒影片)：單次 $5,760 ~ $13,824 TWD (需以兩萬倍數購買)。"
        ]
    )

    # Add images to slides
    import os
    if os.path.exists('asset/img/GoogleVeo1.png') and os.path.exists('asset/img/Kling1.png'):
        slide = prs.slides.add_slide(blank_slide_layout)
        img1 = slide.shapes.add_picture('asset/img/GoogleVeo1.png', Inches(0.5), Inches(1), width=Inches(4))
        img2 = slide.shapes.add_picture('asset/img/Kling1.png', Inches(5), Inches(1), width=Inches(4))
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(1))
        txBox.text_frame.text = "Google Veo vs Kling 模型介面截圖"

    prs.save('d:\\RunWebApp\\KEfilm_Presentation.pptx')
    print("Saved d:\\RunWebApp\\KEfilm_Presentation.pptx")

create_presentation()
